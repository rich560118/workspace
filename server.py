from flask import Flask, Response, request, jsonify, send_file
from flask_cors import CORS
import json
import os
from io import BytesIO
from xml.sax.saxutils import escape
from datetime import date, datetime
import gspread
from google.oauth2.service_account import Credentials

app = Flask(__name__, static_folder='public', static_url_path='')
CORS(app)  # Enable CORS for all routes

# Data file path
DATA_FILE = os.environ.get('DATA_FILE', 'db.json')
DATABASE_URL = os.environ.get('DATABASE_URL')
DATA_STORE_INITIALIZED = False

# Google Sheets configuration
GOOGLE_SHEETS_CREDENTIALS_FILE = 'google-sheets-credentials.json'
GOOGLE_SHEETS_SPREADSHEET_ID = None  # Will be set from environment variable or config
GOOGLE_SHEETS_WORKSHEET_NAME = 'Progress Reports'
VALID_STATUSES = {'未開始', '執行中', '結案'}
PPTX_MIMETYPE = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

def default_data():
    return {
        "workItems": [],
        "progressReports": []
    }

def read_file_data():
    with open(DATA_FILE, 'r', encoding='utf-8') as f:
        return json.load(f)

def write_file_data(data):
    with open(DATA_FILE, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def init_data_file():
    data_dir = os.path.dirname(DATA_FILE)
    if data_dir:
        os.makedirs(data_dir, exist_ok=True)

    if not os.path.exists(DATA_FILE):
        write_file_data(default_data())

def get_db_connection():
    import psycopg

    return psycopg.connect(DATABASE_URL)

def init_postgres_data():
    from psycopg.types.json import Jsonb

    seed_data = read_file_data() if os.path.exists(DATA_FILE) else default_data()

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                CREATE TABLE IF NOT EXISTS app_state (
                    key TEXT PRIMARY KEY,
                    value JSONB NOT NULL,
                    updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
                )
            """)
            cur.execute("""
                INSERT INTO app_state (key, value)
                VALUES (%s, %s)
                ON CONFLICT (key) DO NOTHING
            """, ("data", Jsonb(seed_data)))

def init_data_store():
    global DATA_STORE_INITIALIZED
    if DATA_STORE_INITIALIZED:
        return

    if DATABASE_URL:
        init_postgres_data()
    else:
        init_data_file()

    DATA_STORE_INITIALIZED = True

# Load data from file
def load_data():
    init_data_store()

    if DATABASE_URL:
        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("SELECT value FROM app_state WHERE key = %s", ("data",))
                row = cur.fetchone()
                return row[0] if row else default_data()

    return read_file_data()

# Save data to file
def save_data(data):
    init_data_store()

    if DATABASE_URL:
        from psycopg.types.json import Jsonb

        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO app_state (key, value, updated_at)
                    VALUES (%s, %s, NOW())
                    ON CONFLICT (key) DO UPDATE
                    SET value = EXCLUDED.value,
                        updated_at = NOW()
                """, ("data", Jsonb(data)))
        return

    write_file_data(data)

def next_id(items):
    return max((item.get("id", 0) for item in items), default=0) + 1

def parse_iso_date(value):
    try:
        return date.fromisoformat(value)
    except (TypeError, ValueError):
        return None

def parse_report_month(value):
    try:
        datetime.strptime(value, "%Y-%m")
        return value
    except (TypeError, ValueError):
        return None

def validate_work_item(payload):
    name = (payload.get("name") or "").strip()
    project_member = (payload.get("projectMember") or "").strip()
    start_date = parse_iso_date(payload.get("startDate"))
    end_date = parse_iso_date(payload.get("endDate"))

    if not name:
        return "工作項目名稱不可空白"
    if not start_date or not end_date:
        return "請提供有效的開始日期與結束日期"
    if start_date > end_date:
        return "開始日期不可晚於結束日期"

    payload["name"] = name
    payload["projectMember"] = project_member
    return None

def calculate_progress(work_item, latest_report):
    if not latest_report:
        return 0

    status = latest_report.get("status")
    if status == "結案":
        return 100
    if status == "未開始":
        return 0
    if status != "執行中":
        return 0

    start_date = parse_iso_date(work_item.get("startDate"))
    end_date = parse_iso_date(work_item.get("endDate"))
    execution_date = parse_iso_date(latest_report.get("executionDate")) or date.today()

    if not start_date or not end_date:
        return 0
    if end_date <= start_date:
        return 100 if execution_date >= end_date else 0

    total_days = (end_date - start_date).days + 1
    elapsed_days = (execution_date - start_date).days + 1
    return max(0, min(round((elapsed_days / total_days) * 100, 1), 100))

def as_number(value):
    try:
        return float(value)
    except (TypeError, ValueError):
        return 0

def get_report_month(report):
    if not report:
        return ""

    report_month = parse_report_month(report.get("reportMonth"))
    if report_month:
        return report_month

    execution_date = report.get("executionDate")
    if parse_iso_date(execution_date):
        return execution_date[:7]

    return ""

def format_display_date(value):
    parsed_date = parse_iso_date(value)
    return parsed_date.strftime("%Y/%m/%d") if parsed_date else "未設定"

def format_display_month(value):
    report_month = parse_report_month(value)
    return report_month.replace("-", "/") if report_month else "未設定"

def format_hours(value):
    number = as_number(value)
    return str(int(number)) if number.is_integer() else f"{number:.1f}"

def prepare_progress_report(payload, data):
    report = dict(payload or {})

    try:
        work_item_id = int(report.get("workItemId"))
    except (TypeError, ValueError):
        return None, None, "找不到指定的工作項目"

    work_item = next((item for item in data["workItems"] if item.get("id") == work_item_id), None)
    if not work_item:
        return None, None, "找不到指定的工作項目"

    try:
        actual_hours = float(report.get("actualHours", 0))
    except (TypeError, ValueError):
        return None, None, "實際執行時間格式不正確"

    if actual_hours < 0:
        return None, None, "實際執行時間不可小於 0"

    report_month = parse_report_month(report.get("reportMonth"))
    if not report_month:
        return None, None, "請提供有效的填報月份"
    if not parse_iso_date(report.get("executionDate")):
        return None, None, "請提供有效的實際執行日期"
    if report.get("status") not in VALID_STATUSES:
        return None, None, "請選擇有效的執行情形"

    description = (report.get("description") or "").strip()
    if not description:
        return None, None, "進度說明不可空白"

    return {
        "workItemId": work_item_id,
        "reportMonth": report_month,
        "actualHours": actual_hours,
        "executionDate": report.get("executionDate"),
        "status": report.get("status"),
        "description": description
    }, work_item, None

def report_recency_key(report):
    return (
        get_report_month(report),
        str(report.get("updatedAt") or report.get("createdAt", ""))
    )

def latest_reports_by_work_item(data):
    latest_reports = {}
    for report in data["progressReports"]:
        work_item_id = report.get("workItemId")
        existing = latest_reports.get(work_item_id)
        if not existing or report_recency_key(report) > report_recency_key(existing):
            latest_reports[work_item_id] = report
    return latest_reports

def progress_export_rows(data, report_month=None):
    reports = [
        report for report in data["progressReports"]
        if not report_month or get_report_month(report) == report_month
    ]
    latest_reports = {}
    for report in reports:
        work_item_id = report.get("workItemId")
        existing = latest_reports.get(work_item_id)
        if not existing or report_recency_key(report) > report_recency_key(existing):
            latest_reports[work_item_id] = report

    rows = []

    for work_item in sorted(data["workItems"], key=lambda item: item.get("startDate", "")):
        report = latest_reports.get(work_item.get("id"))
        total_hours = sum(
            as_number(item.get("actualHours", 0))
            for item in reports
            if item.get("workItemId") == work_item.get("id")
        )
        rows.append({
            "name": work_item.get("name", ""),
            "projectMember": work_item.get("projectMember", ""),
            "reportMonth": format_display_month(get_report_month(report)) if report else format_display_month(report_month) if report_month else "尚未填報",
            "actualHours": f"{format_hours(total_hours)} 小時",
            "executionDate": format_display_date(report.get("executionDate")) if report else "尚未填報",
            "status": report.get("status", "未開始") if report else "未開始",
            "description": report.get("description", "") if report else ""
        })

    return rows

def svg_text(text):
    return escape(str(text or ""))

def download_response(body, filename, mimetype):
    return Response(
        body,
        mimetype=mimetype,
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

def build_progress_svg(data, report_month=None):
    rows = progress_export_rows(data, report_month)
    width = 1400
    row_height = 126
    height = 150 + max(len(rows), 1) * row_height
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">',
        '<rect width="100%" height="100%" fill="#f4f6f8"/>',
        '<text x="40" y="60" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="34" font-weight="700" fill="#202a35">進度報告列表</text>',
        f'<text x="40" y="96" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="16" fill="#667789">匯出時間：{datetime.now().strftime("%Y/%m/%d %H:%M")}</text>',
    ]

    if not rows:
        parts.extend([
            '<rect x="40" y="130" width="1320" height="100" rx="10" fill="#ffffff" stroke="#dbe3ea"/>',
            '<text x="700" y="188" text-anchor="middle" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="22" fill="#667789">目前沒有進度報告</text>',
        ])
    else:
        for index, row in enumerate(rows):
            y = 130 + index * row_height
            status_fill = "#e0f2e8" if row["status"] == "結案" else "#dbeef5" if row["status"] == "執行中" else "#fff3d6"
            status_text = "#2f855a" if row["status"] == "結案" else "#19556a" if row["status"] == "執行中" else "#b7791f"
            description = row["description"][:64] + ("..." if len(row["description"]) > 64 else "")

            parts.extend([
                f'<rect x="40" y="{y}" width="1320" height="104" rx="10" fill="#ffffff" stroke="#dbe3ea"/>',
                f'<text x="70" y="{y + 34}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="22" font-weight="700" fill="#202a35">{svg_text(row["name"])}</text>',
                f'<rect x="1200" y="{y + 18}" width="110" height="30" rx="15" fill="{status_fill}"/>',
                f'<text x="1255" y="{y + 39}" text-anchor="middle" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="15" font-weight="700" fill="{status_text}">{svg_text(row["status"])}</text>',
                f'<text x="70" y="{y + 66}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="15" fill="#667789">專案人員</text>',
                f'<text x="70" y="{y + 90}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="18" font-weight="700" fill="#202a35">{svg_text(row["projectMember"] or "未設定")}</text>',
                f'<text x="250" y="{y + 66}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="15" fill="#667789">填報月份</text>',
                f'<text x="250" y="{y + 90}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="18" font-weight="700" fill="#202a35">{svg_text(row["reportMonth"])}</text>',
                f'<text x="430" y="{y + 66}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="15" fill="#667789">累計時數</text>',
                f'<text x="430" y="{y + 90}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="18" font-weight="700" fill="#202a35">{svg_text(row["actualHours"])}</text>',
                f'<text x="610" y="{y + 66}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="15" fill="#667789">執行日期</text>',
                f'<text x="610" y="{y + 90}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="18" font-weight="700" fill="#202a35">{svg_text(row["executionDate"])}</text>',
                f'<text x="800" y="{y + 66}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="15" fill="#667789">進度說明</text>',
                f'<text x="800" y="{y + 90}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="17" fill="#202a35">{svg_text(description or "未填寫")}</text>',
            ])

    parts.append('</svg>')
    return "\n".join(parts)

def build_gantt_export_data(data):
    gantt_data = []
    latest_reports = latest_reports_by_work_item(data)

    for work_item in data["workItems"]:
        related_reports = [
            report for report in data["progressReports"]
            if report.get("workItemId") == work_item.get("id")
        ]
        latest_report = latest_reports.get(work_item.get("id"))
        gantt_data.append({
            "name": work_item.get("name", ""),
            "projectMember": work_item.get("projectMember", ""),
            "start": work_item.get("startDate", ""),
            "end": work_item.get("endDate", ""),
            "progress": calculate_progress(work_item, latest_report),
            "actualHours": sum(as_number(report.get("actualHours", 0)) for report in related_reports),
            "status": latest_report.get("status") if latest_report else "未開始",
        })

    return [
        item for item in gantt_data
        if parse_iso_date(item["start"]) and parse_iso_date(item["end"])
    ]

def build_gantt_svg(data):
    items = sorted(build_gantt_export_data(data), key=lambda item: item["start"])
    width = 1400
    row_height = 74
    height = 150 + max(len(items), 1) * row_height
    label_width = 310
    timeline_x = 350
    timeline_width = 980
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">',
        '<rect width="100%" height="100%" fill="#f4f6f8"/>',
        '<text x="40" y="60" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="34" font-weight="700" fill="#202a35">專案甘特圖</text>',
        f'<text x="40" y="96" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="16" fill="#667789">匯出時間：{datetime.now().strftime("%Y/%m/%d %H:%M")}</text>',
    ]

    if not items:
        parts.extend([
            '<rect x="40" y="130" width="1320" height="100" rx="10" fill="#ffffff" stroke="#dbe3ea"/>',
            '<text x="700" y="188" text-anchor="middle" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="22" fill="#667789">目前沒有工作項目</text>',
            '</svg>'
        ])
        return "\n".join(parts)

    start = min(parse_iso_date(item["start"]) for item in items)
    end = max(parse_iso_date(item["end"]) for item in items)
    total_days = max((end - start).days + 1, 1)

    parts.extend([
        f'<rect x="40" y="125" width="1320" height="{height - 155}" rx="10" fill="#ffffff" stroke="#dbe3ea"/>',
        f'<line x1="{timeline_x}" y1="125" x2="{timeline_x}" y2="{height - 30}" stroke="#dbe3ea"/>',
    ])

    for tick_index in range(0, 6):
        x = timeline_x + (timeline_width * tick_index / 5)
        tick_date = start + (end - start) * tick_index / 5
        parts.extend([
            f'<line x1="{x:.1f}" y1="125" x2="{x:.1f}" y2="{height - 30}" stroke="#e6edf2"/>',
            f'<text x="{x + 4:.1f}" y="148" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="13" fill="#667789">{tick_date.strftime("%Y/%m/%d")}</text>',
        ])

    for index, item in enumerate(items):
        y = 170 + index * row_height
        item_start = parse_iso_date(item["start"])
        item_end = parse_iso_date(item["end"])
        left = timeline_x + ((item_start - start).days / total_days) * timeline_width
        right = timeline_x + (((item_end - start).days + 1) / total_days) * timeline_width
        bar_width = max(right - left, 12)
        progress_width = max(bar_width * item["progress"] / 100, 0)
        label = f'{item["name"]} - {item["projectMember"]}' if item["projectMember"] else item["name"]

        parts.extend([
            f'<line x1="40" y1="{y - 28}" x2="1360" y2="{y - 28}" stroke="#eef2f5"/>',
            f'<text x="70" y="{y}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="18" font-weight="700" fill="#202a35">{svg_text(label)}</text>',
            f'<text x="70" y="{y + 24}" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="13" fill="#667789">{format_display_date(item["start"])} - {format_display_date(item["end"])} · {svg_text(item["status"])} · {format_hours(item["actualHours"])} 小時</text>',
            f'<rect x="{left:.1f}" y="{y - 20}" width="{bar_width:.1f}" height="28" rx="14" fill="#dbeef5" stroke="#9bc7d7"/>',
            f'<rect x="{left:.1f}" y="{y - 20}" width="{progress_width:.1f}" height="28" rx="14" fill="#2f8f83"/>',
            f'<text x="{left + bar_width / 2:.1f}" y="{y - 1}" text-anchor="middle" font-family="Segoe UI, Noto Sans TC, Arial, sans-serif" font-size="13" font-weight="700" fill="#183642">{int(round(item["progress"]))}%</text>',
        ])

    parts.append('</svg>')
    return "\n".join(parts)

def set_run_font(run, size=14, bold=False, color=None):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)

def add_ppt_text(slide, text, left, top, width, height, size=14, bold=False, color=(32, 42, 53)):
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = ""
    run = paragraph.add_run()
    run.text = str(text or "")
    set_run_font(run, size=size, bold=bold, color=color)
    return box

def ppt_bytes(prs):
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

def build_progress_pptx(data, report_month=None):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rows = progress_export_rows(data, report_month)

    add_ppt_text(slide, "進度報告列表", 0.45, 0.3, 8.5, 0.45, size=26, bold=True)
    add_ppt_text(slide, f'匯出時間：{datetime.now().strftime("%Y/%m/%d %H:%M")}', 0.45, 0.78, 5, 0.3, size=10, color=(102, 119, 137))

    headers = ["工作項目", "專案人員", "填報月份", "累計時數", "狀態", "進度說明"]
    widths = [2.2, 1.3, 1.2, 1.1, 1.0, 5.0]
    x = 0.45
    y = 1.25
    for header, width in zip(headers, widths):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.38))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(37, 111, 135)
        shape.line.color.rgb = RGBColor(37, 111, 135)
        add_ppt_text(slide, header, x + 0.05, y + 0.08, width - 0.1, 0.25, size=10, bold=True, color=(255, 255, 255))
        x += width

    for row_index, row in enumerate(rows[:10]):
        y = 1.68 + row_index * 0.48
        values = [row["name"], row["projectMember"], row["reportMonth"], row["actualHours"], row["status"], row["description"][:52]]
        x = 0.45
        for value, width in zip(values, widths):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.45))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(248, 250, 251)
            shape.line.color.rgb = RGBColor(219, 227, 234)
            add_ppt_text(slide, value, x + 0.05, y + 0.08, width - 0.1, 0.28, size=9)
            x += width

    if len(rows) > 10:
        add_ppt_text(slide, f"另有 {len(rows) - 10} 筆資料未顯示於此頁", 0.45, 6.85, 6, 0.3, size=10, color=(102, 119, 137))

    return ppt_bytes(prs)

def build_gantt_pptx(data):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    items = sorted(build_gantt_export_data(data), key=lambda item: item["start"])

    add_ppt_text(slide, "專案甘特圖", 0.45, 0.3, 8.5, 0.45, size=26, bold=True)
    add_ppt_text(slide, f'匯出時間：{datetime.now().strftime("%Y/%m/%d %H:%M")}', 0.45, 0.78, 5, 0.3, size=10, color=(102, 119, 137))

    if not items:
        add_ppt_text(slide, "目前沒有工作項目", 0.45, 1.6, 8, 0.5, size=18, color=(102, 119, 137))
        return ppt_bytes(prs)

    start = min(parse_iso_date(item["start"]) for item in items)
    end = max(parse_iso_date(item["end"]) for item in items)
    total_days = max((end - start).days + 1, 1)
    timeline_left = 3.45
    timeline_width = 8.7
    top = 1.35

    add_ppt_text(slide, "工作項目", 0.55, 1.18, 2.5, 0.3, size=11, bold=True, color=(102, 119, 137))
    add_ppt_text(slide, f"{start.strftime('%Y/%m/%d')} - {end.strftime('%Y/%m/%d')}", timeline_left, 1.18, 4, 0.3, size=11, bold=True, color=(102, 119, 137))

    for index, item in enumerate(items[:8]):
        y = top + index * 0.58
        item_start = parse_iso_date(item["start"])
        item_end = parse_iso_date(item["end"])
        left = timeline_left + ((item_start - start).days / total_days) * timeline_width
        right = timeline_left + (((item_end - start).days + 1) / total_days) * timeline_width
        width = max(right - left, 0.15)
        progress_width = max(width * item["progress"] / 100, 0.03 if item["progress"] > 0 else 0)
        label = f'{item["name"]} - {item["projectMember"]}' if item["projectMember"] else item["name"]

        add_ppt_text(slide, label, 0.55, y + 0.03, 2.65, 0.24, size=10, bold=True)
        add_ppt_text(slide, f'{item["status"]} · {format_hours(item["actualHours"])} 小時', 0.55, y + 0.28, 2.45, 0.22, size=8, color=(102, 119, 137))

        base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(y + 0.12), Inches(width), Inches(0.26))
        base.fill.solid()
        base.fill.fore_color.rgb = RGBColor(219, 238, 245)
        base.line.color.rgb = RGBColor(155, 199, 215)

        if progress_width > 0:
            progress_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(y + 0.12), Inches(progress_width), Inches(0.26))
            progress_shape.fill.solid()
            progress_shape.fill.fore_color.rgb = RGBColor(47, 143, 131)
            progress_shape.line.color.rgb = RGBColor(47, 143, 131)

        add_ppt_text(slide, f'{int(round(item["progress"]))}%', left + width + 0.08, y + 0.13, 0.7, 0.22, size=8, bold=True)

    if len(items) > 8:
        add_ppt_text(slide, f"另有 {len(items) - 8} 個工作項目未顯示於此頁", 0.55, 6.7, 6, 0.3, size=10, color=(102, 119, 137))

    return ppt_bytes(prs)

# Serve the frontend index.html
@app.route('/')
def serve_index():
    return app.send_static_file('index.html')

@app.route('/healthz')
def health_check():
    return jsonify({"status": "ok"})

# Initialize Google Sheets client
def init_google_sheets():
    try:
        if not os.path.exists(GOOGLE_SHEETS_CREDENTIALS_FILE):
            print(f"Google Sheets credentials file not found: {GOOGLE_SHEETS_CREDENTIALS_FILE}")
            return None
        
        # Define the scope
        scope = ['https://spreadsheets.google.com/feeds',
                 'https://www.googleapis.com/auth/drive']
        
        # Load credentials
        creds = Credentials.from_service_account_file(
            GOOGLE_SHEETS_CREDENTIALS_FILE, 
            scopes=scope
        )
        
        # Authorize the client
        client = gspread.authorize(creds)
        return client
    except Exception as e:
        print(f"Error initializing Google Sheets: {e}")
        return None

# Send progress report to Google Sheets
def send_to_google_sheets(progress_report, work_item_name=None, project_member=None):
    try:
        client = init_google_sheets()
        if not client:
            print("Could not initialize Google Sheets client")
            return False
        
        # Get spreadsheet ID from environment or use default
        spreadsheet_id = os.environ.get('GOOGLE_SHEETS_SPREADSHEET_ID', GOOGLE_SHEETS_SPREADSHEET_ID)
        if not spreadsheet_id:
            print("Google Sheets spreadsheet ID not configured")
            return False
        
        # Open the spreadsheet
        spreadsheet = client.open_by_key(spreadsheet_id)
        
        # Try to get the worksheet, create if it doesn't exist
        try:
            worksheet = spreadsheet.worksheet(GOOGLE_SHEETS_WORKSHEET_NAME)
        except gspread.WorksheetNotFound:
            worksheet = spreadsheet.add_worksheet(title=GOOGLE_SHEETS_WORKSHEET_NAME, rows=1000, cols=20)
            # Add headers
            worksheet.append_row(['Timestamp', 'Work Item ID', 'Work Item Name', 'Project Member', 'Report Month', 'Actual Hours', 'Execution Date', 'Status', 'Description', 'Report ID'])
        
        # Prepare row data
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        row_data = [
            timestamp,
            progress_report.get('workItemId', ''),
            work_item_name or '',
            project_member or '',
            progress_report.get('reportMonth', ''),
            progress_report.get('actualHours', ''),
            progress_report.get('executionDate', ''),
            progress_report.get('status', ''),
            progress_report.get('description', ''),
            progress_report.get('id', '')
        ]
        
        # Append row to worksheet
        worksheet.append_row(row_data)
        return True
    except Exception as e:
        print(f"Error sending data to Google Sheets: {e}")
        return False

# API Routes

# Get all work items
@app.route('/api/work-items', methods=['GET'])
def get_work_items():
    data = load_data()
    return jsonify(data["workItems"])

# Create a new work item
@app.route('/api/work-items', methods=['POST'])
def create_work_item():
    data = load_data()
    new_item = request.get_json(silent=True) or {}
    validation_error = validate_work_item(new_item)
    if validation_error:
        return jsonify({"error": validation_error}), 400
    
    # Add ID and timestamps
    new_item["id"] = next_id(data["workItems"])
    new_item["createdAt"] = datetime.now().isoformat()
    new_item["updatedAt"] = datetime.now().isoformat()

    data["workItems"].append(new_item)
    save_data(data)
    
    return jsonify(new_item), 201

# Update a work item
@app.route('/api/work-items/<int:item_id>', methods=['PUT'])
def update_work_item(item_id):
    data = load_data()
    update_data = request.get_json(silent=True) or {}
    
    for item in data["workItems"]:
        if item["id"] == item_id:
            candidate = {**item, **update_data}
            validation_error = validate_work_item(candidate)
            if validation_error:
                return jsonify({"error": validation_error}), 400

            item.update(update_data)
            item["name"] = candidate["name"]
            item["projectMember"] = candidate.get("projectMember", "")
            item["updatedAt"] = datetime.now().isoformat()

            save_data(data)
            return jsonify(item)
    
    return jsonify({"error": "Work item not found"}), 404

# Delete a work item
@app.route('/api/work-items/<int:item_id>', methods=['DELETE'])
def delete_work_item(item_id):
    data = load_data()
    
    for i, item in enumerate(data["workItems"]):
        if item["id"] == item_id:
            deleted_item = data["workItems"].pop(i)
            data["progressReports"] = [
                report for report in data["progressReports"]
                if report.get("workItemId") != item_id
            ]
            save_data(data)
            return jsonify(deleted_item)
    
    return jsonify({"error": "Work item not found"}), 404

# Get all progress reports
@app.route('/api/progress-reports', methods=['GET'])
def get_progress_reports():
    data = load_data()
    return jsonify(data["progressReports"])

# Create a new progress report
@app.route('/api/progress-reports', methods=['POST'])
def create_progress_report():
    data = load_data()
    new_report, work_item, validation_error = prepare_progress_report(request.get_json(silent=True), data)
    if validation_error:
        return jsonify({"error": validation_error}), 400

    existing_report = next((
        report for report in sorted(data["progressReports"], key=report_recency_key, reverse=True)
        if report.get("workItemId") == new_report["workItemId"]
        and get_report_month(report) == new_report["reportMonth"]
    ), None)

    if existing_report:
        existing_report.update(new_report)
        existing_report["updatedAt"] = datetime.now().isoformat()
        save_data(data)
        return jsonify(existing_report)
    
    # Add ID and timestamp
    new_report["id"] = next_id(data["progressReports"])
    new_report["createdAt"] = datetime.now().isoformat()
    
    data["progressReports"].append(new_report)
    save_data(data)
    
    # Also send to Google Sheets if configured
    try:
        # Get the work item name for better context in Google Sheets
        # Send to Google Sheets (non-blocking - don't fail the request if this fails)
        send_to_google_sheets(new_report, work_item.get("name"), work_item.get("projectMember"))
    except Exception as e:
        print(f"Error in Google Sheets integration: {e}")
        # Continue anyway - we don't want to fail the main request
    
    return jsonify(new_report), 201

# Update an existing progress report
@app.route('/api/progress-reports/<int:report_id>', methods=['PUT'])
def update_progress_report(report_id):
    data = load_data()
    updates, _work_item, validation_error = prepare_progress_report(request.get_json(silent=True), data)
    if validation_error:
        return jsonify({"error": validation_error}), 400

    for report in data["progressReports"]:
        if report.get("id") == report_id:
            report.update(updates)
            report["updatedAt"] = datetime.now().isoformat()
            save_data(data)
            return jsonify(report)

    return jsonify({"error": "找不到指定的進度報告"}), 404

# Get progress reports for a specific work item
@app.route('/api/progress-reports/work-item/<int:work_item_id>', methods=['GET'])
def get_progress_reports_by_work_item(work_item_id):
    data = load_data()
    reports = [report for report in data["progressReports"] if report.get("workItemId") == work_item_id]
    return jsonify(reports)

# Get Gantt chart data
@app.route('/api/gantt-data', methods=['GET'])
def get_gantt_data():
    data = load_data()
    gantt_data = []
    
    for work_item in data["workItems"]:
        total_actual_hours = 0
        related_reports = []
        for report in data["progressReports"]:
            if report.get("workItemId") == work_item.get("id"):
                total_actual_hours += as_number(report.get("actualHours", 0))
                related_reports.append(report)

        latest_report = None
        if related_reports:
            latest_report = max(
                related_reports,
                key=report_recency_key
            )
        
        gantt_data.append({
            "id": work_item.get("id"),
            "name": work_item.get("name", ""),
            "projectMember": work_item.get("projectMember", ""),
            "start": work_item.get("startDate", ""),
            "end": work_item.get("endDate", ""),
            "progress": calculate_progress(work_item, latest_report),
            "actualHours": total_actual_hours,
            "latestStatus": latest_report.get("status") if latest_report else "未開始",
            "latestReportMonth": get_report_month(latest_report),
            "latestExecutionDate": latest_report.get("executionDate") if latest_report else "",
            "latestDescription": latest_report.get("description") if latest_report else ""
        })
    
    return jsonify(gantt_data)

@app.route('/api/export/progress-reports.svg', methods=['GET'])
def export_progress_reports_svg():
    data = load_data()
    report_month = request.args.get("reportMonth")
    if report_month and not parse_report_month(report_month):
        return jsonify({"error": "請提供有效的填報月份"}), 400
    return download_response(build_progress_svg(data, report_month), "progress-reports.svg", "image/svg+xml")

@app.route('/api/export/gantt.svg', methods=['GET'])
def export_gantt_svg():
    data = load_data()
    return download_response(build_gantt_svg(data), "gantt-chart.svg", "image/svg+xml")

@app.route('/api/export/progress-reports.pptx', methods=['GET'])
def export_progress_reports_pptx():
    data = load_data()
    report_month = request.args.get("reportMonth")
    if report_month and not parse_report_month(report_month):
        return jsonify({"error": "請提供有效的填報月份"}), 400
    return send_file(
        build_progress_pptx(data, report_month),
        as_attachment=True,
        download_name="progress-reports.pptx",
        mimetype=PPTX_MIMETYPE
    )

@app.route('/api/export/gantt.pptx', methods=['GET'])
def export_gantt_pptx():
    data = load_data()
    return send_file(
        build_gantt_pptx(data),
        as_attachment=True,
        download_name="gantt-chart.pptx",
        mimetype=PPTX_MIMETYPE
    )

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    debug = os.environ.get('FLASK_DEBUG') == '1'
    app.run(host='0.0.0.0', debug=debug, port=port)
